import fitz  # PyMuPDF
import docx
from pptx import Presentation

def extract_text_from_pdf(path):
    doc = fitz.open(path)
    text = "\n".join(page.get_text().strip() for page in doc)
    return "\n".join([line for line in text.splitlines() if line.strip()])

def extract_text_from_docx(path):
    doc = docx.Document(path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)

def extract_text_from_pptx(path):
    prs = Presentation(path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)

def extract_text_from_txt(path):
    try:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(path, "r", encoding="latin-1") as f:
            return f.read()
