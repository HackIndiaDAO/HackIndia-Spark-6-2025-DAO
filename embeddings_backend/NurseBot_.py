# NurseBot 
import docx, faiss, numpy as np, os, google.generativeai as genai
from sentence_transformers import SentenceTransformer
from PyPDF2 import PdfReader
import pandas as pd
from pptx import Presentation
from dotenv import load_dotenv

# Load environment
load_dotenv()
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"
genai.configure(api_key=os.getenv("API_KEY"))
gemini = genai.GenerativeModel("gemini-1.5-flash")

embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
print("Model loaded successfully!")

# Loaders
from docx import Document

def load_docx(path):
    try:
        doc = Document(path)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    except Exception as e:
        print(f"[ERROR] Failed to read DOCX file: {e}")
        return ""

def load_pdf(path):
    try:
        reader = PdfReader(path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        print(f"[ERROR] Failed to read PDF file: {e}")
        return ""
    
def load_txt(path):
    try:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        print(f"[ERROR] Failed to read TXT file: {e}")
        return ""
    
def load_pptx(path):
    try:
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"[ERROR] Failed to read PPTX file: {e}")
        return ""
    

def load_excel(path):
    try:
        dfs = pd.read_excel(path, sheet_name=None)  # Load all sheets
        text = ""
        for name, df in dfs.items():
            text += f"--- Sheet: {name} ---\n"
            text += df.to_string(index=False)
            text += "\n\n"
        return text
    except Exception as e:
        print(f"[ERROR] Failed to read Excel file: {e}")
        return ""
    
    
# Chunker
def chunk_text(text, chunk_size=500):
    paragraphs = text.split("\n")
    chunks = []
    current_chunk = ""
    page = 1

    for para in paragraphs:
        if len(current_chunk) + len(para) + 1 <= chunk_size:
            current_chunk += para + "\n"
        else:
            chunks.append((page, current_chunk.strip()))
            page += 1
            current_chunk = para + "\n"

    if current_chunk.strip():
        chunks.append((page, current_chunk.strip()))

    return chunks


# Indexing
def build_index(chunks):
    # Extract embeddings for each chunk
    embeddings = embedding_model.encode([chunk[1] for chunk in chunks], convert_to_numpy=True)
    print("HELPPPPPPPP")
    print(embeddings)
    # Create a FAISS index
    dim = embeddings.shape[1]
    index = faiss.IndexFlatL2(dim)
    index.add(embeddings)
    
    # Return the embeddings along with chunks
    return index, chunks, embeddings

# Loader dispatcher
def load_file(file_path):
    print(f"[DEBUG] Loading file: {file_path}")
    if file_path.endswith(".docx"):
        return load_docx(file_path)
    elif file_path.endswith(".pdf"):
        return load_pdf(file_path)
    elif file_path.endswith(".txt"):
        return load_txt(file_path)
    elif file_path.endswith(".pptx"):
        return load_pptx(file_path)
    elif file_path.endswith(".xlsx"):
        return load_excel(file_path)
    else:
        raise ValueError("Unsupported file type")

# QA Function
def ask_question(query, index, chunks, k=3):
    query_vec = embedding_model.encode([query], convert_to_numpy=True)
    _, I = index.search(query_vec, k)
    context = "\n".join([f"Page {chunks[i][0]}: {chunks[i][1]}" for i in I[0]])

    prompt = f"""
You are a helpful assistant. Use the following document content to answer the user's question. Context:{context}
Answer the way it is in the document uploaded.
If possible, mention the paragraph or the page number it is on.
Understand the user based on their natural language.
...
Question: {query}
"""
    response = gemini.generate_content(prompt)
    return response.text if response.text else "Sorry, I couldn't generate an answer from the document."